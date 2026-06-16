# -*- coding: utf-8 -*-
"""Generate 2-slide PPT: CUHK housing gap — concise Colliers vs calibrated."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT_PATH = ROOT / "學生宿舍" / "中大住宿缺口测算.pptx"

FONT_NAME = "Microsoft YaHei"
CHAPTER = "二、學生住宿市場"
FONT_CHAPTER = 24
FONT_SECTION = 18
FONT_NARRATIVE = 13
FONT_TABLE = 10
FONT_TABLE_SM = 10
FONT_CONCLUSION = 12

CSC_RED = RGBColor(192, 0, 0)
TEXT_DARK = RGBColor(51, 51, 51)
TEXT_GRAY = RGBColor(102, 102, 102)
WHITE = RGBColor(255, 255, 255)
ROW_ALT = RGBColor(248, 248, 248)
BAR_BG = RGBColor(220, 220, 220)
BAR_GRAY = RGBColor(160, 160, 160)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = 0.35
CONTENT_W = 12.63
HALF_W = 6.05
GAP = 0.53
RIGHT_X = MARGIN + HALF_W + GAP

# CUHK Facts 2024（政府资助课程，截至 2024/9/30）
UG_TOTAL, UG_LOCAL, UG_NL = 18438, 15742, 2696
PG_TOTAL, PG_LOCAL, PG_NL = 4650, 1699, 2951
ALL_TOTAL = UG_TOTAL + PG_TOTAL          # 23,088
ALL_LOCAL = UG_LOCAL + PG_LOCAL          # 17,441
ALL_NL = UG_NL + PG_NL                   # 5,647
UG_EFF = round(UG_LOCAL * 0.35 + UG_NL)  # 8,206
PG_EFF = round(PG_LOCAL * 0.35 + PG_NL)  # 3,546
ALL_EFF = UG_EFF + PG_EFF                # 11,752
SUPPLY = 10600
GAP = ALL_EFF - SUPPLY                   # ~1,152

# 内容区纵向分区（英寸）
Y_TITLE = 0.82
Y_NARR = 1.15
H_NARR = 1.05
Y_BODY = 2.35
H_BODY = 3.55
Y_CONC = 6.05


def set_run_font(run, size: int, bold: bool = False, color: RGBColor = TEXT_DARK):
    run.font.name = FONT_NAME
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_rich_paragraph(slide, left, top, width, height, segments, size=FONT_NARRATIVE, line_spacing=1.35):
    """segments: [(text, color, bold), ...]"""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.line_spacing = line_spacing
    p.alignment = PP_ALIGN.LEFT
    for i, (text, color, bold) in enumerate(segments):
        if i == 0:
            p.text = text
            run = p.runs[0]
        else:
            run = p.add_run()
            run.text = text
        set_run_font(run, size, bold, color)
    return box


def add_header(slide):
    ph = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.35), Inches(0.18), Inches(1.6), Inches(0.45)
    )
    ph.fill.background()
    ph.line.fill.background()
    box = slide.shapes.add_textbox(Inches(0.35), Inches(0.22), Inches(1.6), Inches(0.4))
    p = box.text_frame.paragraphs[0]
    p.text = "【Logo】"
    p.alignment = PP_ALIGN.CENTER
    set_run_font(p.runs[0], 10, False, TEXT_GRAY)

    box = slide.shapes.add_textbox(Inches(8.5), Inches(0.15), Inches(4.5), Inches(0.55))
    p = box.text_frame.paragraphs[0]
    p.text = CHAPTER
    p.alignment = PP_ALIGN.RIGHT
    set_run_font(p.runs[0], FONT_CHAPTER, True, CSC_RED)

    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.35), Inches(0.72), Inches(12.6), Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = CSC_RED
    line.line.fill.background()


def add_subsection_title(slide, title: str):
    sq = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.35), Inches(Y_TITLE), Inches(0.14), Inches(0.26)
    )
    sq.fill.solid()
    sq.fill.fore_color.rgb = CSC_RED
    sq.line.fill.background()
    box = slide.shapes.add_textbox(Inches(0.56), Inches(Y_TITLE - 0.03), Inches(12.2), Inches(0.42))
    p = box.text_frame.paragraphs[0]
    p.text = title
    set_run_font(p.runs[0], FONT_SECTION, True, CSC_RED)


def add_table_banner(slide, left, top, width, text: str):
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.36)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = CSC_RED
    bar.line.fill.background()
    tf = bar.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    set_run_font(p.runs[0], FONT_TABLE, True, WHITE)


def style_table_cell(cell, text, bold=False, color=TEXT_DARK, size=FONT_TABLE_SM, align=PP_ALIGN.CENTER, fill=None):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.alignment = align
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Pt(5)
    cell.margin_right = Pt(5)
    cell.margin_top = Pt(4)
    cell.margin_bottom = Pt(4)
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    if p.runs:
        set_run_font(p.runs[0], size, bold, color)


def add_data_table(slide, left, top, width, height, headers, rows, col_widths=None,
                   highlight_rows=None, left_cols=None):
    if left_cols is None:
        left_cols = {0}
    ts = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(left), Inches(top), Inches(width), Inches(height))
    table = ts.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    for j, h in enumerate(headers):
        style_table_cell(table.cell(0, j), h, bold=True, color=WHITE, size=FONT_TABLE, fill=CSC_RED)
    for i, row in enumerate(rows):
        bg = ROW_ALT if i % 2 else WHITE
        hi = highlight_rows and i in highlight_rows
        for j, val in enumerate(row):
            align = PP_ALIGN.LEFT if j in left_cols else PP_ALIGN.CENTER
            c = CSC_RED if hi else TEXT_DARK
            style_table_cell(table.cell(i + 1, j), val, bold=hi, color=c, align=align, fill=bg)
    return ts


def add_hbar_chart(slide, left, top, width, height, title, bars, max_val=None):
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = ROW_ALT
    panel.line.color.rgb = RGBColor(210, 210, 210)

    add_table_banner(slide, left + 0.08, top + 0.1, width - 0.16, title)

    if max_val is None:
        max_val = max(v for _, v, _ in bars) * 1.12
    track_left = left + 0.28
    track_w = width - 1.35
    y = top + 0.62
    row_h = (height - 0.75) / len(bars)

    for label, value, color in bars:
        ratio = min(value / max_val, 1.0) if max_val else 0
        box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(y), Inches(1.85), Inches(row_h * 0.5))
        p = box.text_frame.paragraphs[0]
        p.text = label
        set_run_font(p.runs[0], 10, False, TEXT_DARK)

        track = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(track_left), Inches(y + 0.05), Inches(track_w), Inches(row_h * 0.38),
        )
        track.fill.solid()
        track.fill.fore_color.rgb = BAR_BG
        track.line.fill.background()
        fill = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(track_left), Inches(y + 0.05),
            Inches(max(track_w * ratio, 0.08)), Inches(row_h * 0.38),
        )
        fill.fill.solid()
        fill.fill.fore_color.rgb = color or CSC_RED
        fill.line.fill.background()

        val_text = f"{value:,}" if isinstance(value, int) else str(value)
        box = slide.shapes.add_textbox(
            Inches(track_left + track_w + 0.06), Inches(y), Inches(0.9), Inches(row_h * 0.5)
        )
        p = box.text_frame.paragraphs[0]
        p.text = val_text
        set_run_font(p.runs[0], 10, True, CSC_RED)
        y += row_h


def add_conclusion(slide, text: str):
    box = slide.shapes.add_textbox(Inches(MARGIN), Inches(Y_CONC), Inches(CONTENT_W), Inches(0.55))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.25
    set_run_font(p.runs[0], FONT_CONCLUSION, True, CSC_RED)


def slide_colliers(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide)
    add_subsection_title(slide, "1、中大住宿需求测算（Colliers 口径）")

    add_rich_paragraph(slide, MARGIN, Y_NARR, CONTENT_W, H_NARR, [
        ("▶ 2024/25学年香港中文大学政府资助全日制在校生约", TEXT_DARK, False),
        ("2.31万人", CSC_RED, True),
        ("（本科", TEXT_DARK, False),
        ("1.84万", CSC_RED, True),
        ("、研究生", TEXT_DARK, False),
        ("0.47万", CSC_RED, True),
        ("），其中本地生约", TEXT_DARK, False),
        ("1.74万人", CSC_RED, True),
        ("、非本地生约", TEXT_DARK, False),
        ("0.56万人", CSC_RED, True),
        ("。参照高力国际测算逻辑，须住宿人数＝本地全日制×", TEXT_DARK, False),
        ("35%", CSC_RED, True),
        ("＋非本地×", TEXT_DARK, False),
        ("100%", CSC_RED, True),
        ("，约合", TEXT_DARK, False),
        ("1.18万", CSC_RED, True),
        ("床；校内供给约", TEXT_DARK, False),
        ("1.06万", CSC_RED, True),
        ("床，满足率约", TEXT_DARK, False),
        ("90%", CSC_RED, True),
        ("，净缺口约", TEXT_DARK, False),
        ("1,150", CSC_RED, True),
        ("床。", TEXT_DARK, False),
    ])

    add_table_banner(slide, MARGIN, Y_BODY, HALF_W, "供需测算一览")
    add_data_table(
        slide, MARGIN, Y_BODY + 0.42, HALF_W, H_BODY - 0.5,
        ["群体", "在校生", "有效需求", "校内供给"],
        [
            ["本科生", f"{UG_TOTAL:,}", f"{UG_EFF:,}", "~9,500"],
            ["研究生", f"{PG_TOTAL:,}", f"{PG_EFF:,}", "~1,100"],
            ["合计", f"{ALL_TOTAL:,}", f"{ALL_EFF:,}", f"~{SUPPLY:,}"],
        ],
        col_widths=[1.1, 1.45, 1.45, 1.45],
        highlight_rows={2}, left_cols={0},
    )

    add_hbar_chart(
        slide, RIGHT_X, Y_BODY, HALF_W, H_BODY,
        "需求 vs 供给（床）",
        [
            ("有效需求", ALL_EFF, CSC_RED),
            ("校内供给", SUPPLY, BAR_GRAY),
            ("净缺口", GAP, CSC_RED),
        ],
        max_val=13000,
    )

    add_conclusion(slide, "结论：按 Colliers 行业口径，中大住宿净缺口约 1,150 床，有效需求满足率约 90%。")


def slide_calibrated(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide)
    add_subsection_title(slide, "2、中大住宿缺口测算（政策校准口径）")

    add_rich_paragraph(slide, MARGIN, Y_NARR, CONTENT_W, H_NARR, [
        ("▶ 考虑中大书院制度及研究生政策校准后，本科生校内覆盖率仅约", TEXT_DARK, False),
        ("52%", CSC_RED, True),
        ("，未获宿及高年级搬出形成约", TEXT_DARK, False),
        ("7,500–9,000", CSC_RED, True),
        ("人次/年校外需求；研究生非本地占", TEXT_DARK, False),
        ("63%", CSC_RED, True),
        ("，修课式不提供校内宿，研究式竞", TEXT_DARK, False),
        ("1,100", CSC_RED, True),
        ("床研宿。校准后可寻址校外住宿需求约", TEXT_DARK, False),
        ("1.15–1.35万", CSC_RED, True),
        ("人，远超周边可核实商业化床位（仅数百个）。", TEXT_DARK, False),
    ])

    add_table_banner(slide, MARGIN, Y_BODY, HALF_W, "分群体缺口一览")
    add_data_table(
        slide, MARGIN, Y_BODY + 0.42, HALF_W, H_BODY - 0.5,
        ["群体", "在校生", "校内宿位", "校外需求"],
        [
            ["本科生", f"{UG_TOTAL:,}", "~9,500", "7,500–9,000 人次/年"],
            ["研究生", f"{PG_TOTAL:,}", "~1,100", "~4,000 人头"],
            ["合计", f"{ALL_TOTAL:,}", f"~{SUPPLY:,}", "11,500–13,500"],
        ],
        col_widths=[0.95, 1.2, 1.2, 2.3],
        highlight_rows={2}, left_cols={0, 3},
    )

    add_hbar_chart(
        slide, RIGHT_X, Y_BODY, HALF_W, H_BODY,
        "两口径缺口对比（床）",
        [
            ("Colliers 净缺口", GAP, BAR_GRAY),
            ("校准可寻址需求", 12500, CSC_RED),
        ],
        max_val=16000,
    )

    add_conclusion(slide, "结论：投资测算建议采用校准口径——中大可寻址校外需求约 1.15–1.35 万人，结构性缺口极大。")


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide_colliers(prs)
    slide_calibrated(prs)
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    print(f"Saved: {OUT_PATH}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
