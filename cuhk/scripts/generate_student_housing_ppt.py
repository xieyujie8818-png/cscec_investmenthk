# -*- coding: utf-8 -*-
"""Generate 5-slide PPT: CUHK student housing demand & off-campus apartments (v5)."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT_PATH = ROOT / "中大學生住宿需求與校外公寓.pptx"
OUT_PATH_FALLBACK = ROOT / "中大學生住宿需求與校外公寓_v8.pptx"

# ── 統一版式 ──
FONT_NAME = "Microsoft YaHei"
CHAPTER = "二、學生住宿市場"
FONT_CHAPTER = 24
FONT_SECTION = 18
FONT_BODY = 13
FONT_TABLE = 11
FONT_TABLE_SM = 10
FONT_FOOTER = 9

CSC_RED = RGBColor(192, 0, 0)
TEXT_DARK = RGBColor(51, 51, 51)
TEXT_GRAY = RGBColor(102, 102, 102)
WHITE = RGBColor(255, 255, 255)
ROW_ALT = RGBColor(248, 248, 248)
ROW_PEACH = RGBColor(255, 236, 210)
BAR_BG = RGBColor(220, 220, 220)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = 0.35
CONTENT_W = 12.63
LEASE_MONTHS = 11.5
ACADEMIC_MONTHS = 9

# HK01 2025-07-24 租金；面積為估算（見 slide_parkwood_detail 註腳）
# 按 richitt 批則住宅總樓面 41,385呎÷122戶≈339呎/戶，再依 HK01 開放式/一房/兩房典型戶型推算
PARKWOOD_GFA_TOTAL = 41385
PARKWOOD_UNITS = 122
PARKWOOD_ROOMS = [
    {"label": "四人房（床位）", "monthly": 5335, "unit_area": 400, "occupants": 4},
    {"label": "三人房（床位）", "monthly": 5950, "unit_area": 330, "occupants": 3},
    {"label": "雙人房（床位）", "monthly": 5950, "unit_area": 280, "occupants": 2},
]

# 圖2：學生公寓 → 長租酒店；月租 OSA 2025 + HK01/平台
# 欄位：序号, 类型, 项目, 宿位, 月租（港币）, 运营方, 年份
MARKET_DORMS = [
    ["1", "学生公寓", "迎林 Parkwood", "~360", "5,335–5,950/床", "恒基（中大2025认可）", "2025"],
    ["2", "学生公寓", "安怡居·禾睦居", "21", "约9,500/套", "安怡居（沙田近中大）", "2021"],
    ["3", "学生公寓", "HereIN 中大府邸", "~16", "8,300起/月", "HereIN / 集好家", "2020"],
    ["4", "学生公寓", "酷點 Kudian", "分散床位", "4,800–7,000起/床", "大围/大埔/马鞍山/金禧等", "2020"],
    ["5", "学生公寓", "安怡居（全港共居）", "多区", "3,800起/床", "安怡居有限公司", "2012"],
    ["6", "长租酒店", "Aspire House\n麗豪酒店", "双人床位", "8,750起/床", "富豪酒店集团", "2026"],
    ["7", "长租酒店", "ALVA 帝逸酒店", "合租床位", "8,550起/床", "皇家 / 新鸿基", "2024"],
    ["8", "长租酒店", "Royal Park 帝都", "合租床位", "11,880起/床", "帝都酒店", "2023"],
]

MARKET_DORM_HEADERS = ["序号", "类型", "项目", "宿位", "月租（港币）", "运营方/投资方", "投入运营"]

# 圖1（續）：九书院 + 校内设施合表
CAMPUS_ALL = [
    ["崇基學院", "1,448", "四年一宿·申請制", "16,826/年", "按書院計分分配"],
    ["新亞書院", "1,470", "四年一宿·申請制", "17,280/年", "N-2：首兩年後須遴選"],
    ["聯合書院", "1,538", "四年一宿·申請制", "見書院公布", "非本地生保證首兩年"],
    ["逸夫書院", "1,211", "四年一宿·申請制", "見書院公布", "按書院計分分配"],
    ["晨興書院", "300", "全宿共膳·四年住校", "見書院公布", "非本地生不可選校外"],
    ["善衡書院", "600", "全宿共膳·四年住校", "見書院公布", "非本地生不可選校外"],
    ["敬文書院", "302", "全宿共膳·四年住校", "17,774/年", "非本地生不可選校外"],
    ["伍宜孫書院", "604", "四年一宿·申請制", "見書院公布", "按書院計分分配"],
    ["和聲書院", "640", "四年一宿·申請制", "見書院公布", "按書院計分分配"],
    ["國際生堂 I-House", "~598", "本科生", "OSA 管理", "非书院设施"],
    ["研究生宿舍 PGH", "1,000–1,500", "研究式/教資會修課式", "1,685–4,055/月", "修课式研无校内宿"],
    ["蔡繼有宿舍（聯合）", "250", "2025 啟用", "聯合書院官網", "新建擴容"],
    ["崇基 E 計劃雙子樓", "+300（在建）", "2026/27 落成", "中大新聞稿", "在建擴容"],
]

# 星岛头条 2025：罗湖服务式公寓 / 中大港漂跨境租住
SHENZHEN_RENTAL = [
    ["形成原因", "未获校内宿舍", "授课式硕士尤甚；港漂北上成本替代"],
    ["市场观察", "身边同学跨境上学普遍", "星岛头条采访中大内地生"],
    ["罗湖服务式公寓", "约40%租客为学生", "辉盛名致试业；325套房"],
    ["租客画像", "在港攻读硕士内地生", "近9成选单人Studio"],
    ["香港合租", "1.5–3万元/月", "校内附近400–600呎，2–3人"],
    ["香港学生宿舍", "5,000–1.4万元/月", "商业宿舍80–200呎"],
    ["深圳福田口岸", "一房约5千元/月", "约400+呎；两房约7千元"],
    ["中大跨境通勤", "约40–60分钟", "东铁大学站↔罗湖/福田/莲塘"],
]

UG_DEMAND = [
    ["在校生（2024）", "18,438", "非本地 2,696"],
    ["校内宿位", "~9,500", "九书院 + I-House"],
    ["校内覆盖率", "约 52%", "9,500 / 18,438"],
    ["校外需求（粗算）", "9,000+ 人次", "未获宿 + 高年级搬出"],
]

PG_DEMAND = [
    ["在校生（2024）", "6,349", "非本地 1,699"],
    ["PGH 研宿", "~1,100", "研究式 / 教资会修课式"],
    ["校内覆盖率", "约 17%", "1,100 / 6,349"],
    ["结构性缺口", "修课式/兼读/自资", "OSA：不提供校内住宿"],
]

COLLEGES = CAMPUS_ALL[:9]

FULL_RES_COLLEGES = {"晨興書院", "善衡書院", "敬文書院"}

OFFCAMPUS_CONCLUSION = (
    "结论：中大周边商业化学生宿舍供给严重不足——可核实床位仅约数百个量级，"
    "相对本科校外需求 9,000+ 人次/学年，结构性缺口极大；酒店长租可补充但单价显著更高。"
)


def set_run_font(run, size: int, bold: bool = False, color: RGBColor = TEXT_DARK):
    run.font.name = FONT_NAME
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_textbox(slide, left, top, width, height, lines, size=FONT_BODY, bold=False,
                color=TEXT_DARK, align=PP_ALIGN.LEFT, line_spacing=1.3):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.line_spacing = line_spacing
        if p.runs:
            set_run_font(p.runs[0], size, bold, color)
    return box


def add_header(slide):
    placeholder = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.35), Inches(0.18), Inches(1.6), Inches(0.45)
    )
    placeholder.fill.background()
    placeholder.line.fill.background()
    add_textbox(slide, 0.35, 0.22, 1.6, 0.4, ["【Logo】"], size=10, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
    add_textbox(slide, 8.5, 0.15, 4.5, 0.55, [CHAPTER], size=FONT_CHAPTER, bold=True, color=CSC_RED, align=PP_ALIGN.RIGHT)
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.35), Inches(0.72), Inches(12.6), Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = CSC_RED
    line.line.fill.background()


def add_subsection_title(slide, title: str, top: float = 0.82):
    square = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.35), Inches(top), Inches(0.14), Inches(0.26)
    )
    square.fill.solid()
    square.fill.fore_color.rgb = CSC_RED
    square.line.fill.background()
    add_textbox(slide, 0.56, top - 0.03, 12.2, 0.42, [title], size=FONT_SECTION, bold=True, color=CSC_RED)


def add_table_banner(slide, left, top, width, text: str):
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.38)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = CSC_RED
    bar.line.fill.background()
    tf = bar.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    set_run_font(p.runs[0], FONT_TABLE, True, WHITE)


def add_footer(slide, text: str):
    add_textbox(slide, MARGIN, 7.05, CONTENT_W, 0.35, [text], size=FONT_FOOTER, color=TEXT_GRAY, align=PP_ALIGN.CENTER)


def style_table_cell(cell, text, bold=False, color=TEXT_DARK, size=FONT_TABLE, align=PP_ALIGN.CENTER, fill=None):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.alignment = align
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Pt(4)
    cell.margin_right = Pt(4)
    cell.margin_top = Pt(3)
    cell.margin_bottom = Pt(3)
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    if p.runs:
        set_run_font(p.runs[0], size, bold, color)


def add_data_table(
    slide, left, top, width, height, headers, rows,
    col_widths=None, header_size=FONT_TABLE, body_size=FONT_TABLE_SM,
    highlight_rows=None, peach_rows=False, left_cols=None,
):
    if left_cols is None:
        left_cols = {0, 1}
    table_shape = slide.shapes.add_table(
        len(rows) + 1, len(headers), Inches(left), Inches(top), Inches(width), Inches(height)
    )
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for j, h in enumerate(headers):
        style_table_cell(table.cell(0, j), h, bold=True, color=WHITE, size=header_size, fill=CSC_RED)

    for i, row in enumerate(rows):
        if peach_rows:
            bg = ROW_PEACH if i % 2 else WHITE
        else:
            bg = ROW_ALT if i % 2 else WHITE
        hi = highlight_rows and i in highlight_rows
        for j, val in enumerate(row):
            align = PP_ALIGN.LEFT if j in left_cols else PP_ALIGN.CENTER
            c = CSC_RED if hi and j == 2 else TEXT_DARK
            b = hi and j == 2
            style_table_cell(table.cell(i + 1, j), val, bold=b, color=c, size=body_size, align=align, fill=bg)
    return table_shape


def add_kpi_box(slide, left, top, width, height, label, value, sub=""):
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = ROW_ALT
    panel.line.color.rgb = RGBColor(210, 210, 210)
    add_textbox(slide, left + 0.1, top + 0.08, width - 0.2, 0.28, [label], size=9, color=TEXT_GRAY)
    add_textbox(
        slide, left + 0.1, top + 0.32, width - 0.2, 0.38,
        [value], size=20, bold=True, color=CSC_RED, align=PP_ALIGN.CENTER,
    )
    if sub:
        add_textbox(slide, left + 0.1, top + 0.72, width - 0.2, 0.22, [sub], size=8, color=TEXT_GRAY, align=PP_ALIGN.CENTER)


def add_gap_bars(slide, left, top, width, height, bars=None):
    if bars is None:
        bars = [
            ("本科生校內獲宿", 9500, 18438, "52%"),
            ("研究生校內研宿", 1100, 6349, "17%"),
        ]
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = ROW_ALT
    panel.line.color.rgb = RGBColor(210, 210, 210)
    bar_max_w = width - 1.1
    y = top + 0.12
    for label, num, denom, pct in bars:
        ratio = min(num / denom, 1.0) if denom else 0
        add_textbox(slide, left + 0.12, y, 3.2, 0.28, [label], size=10)
        track = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(left + 3.35), Inches(y + 0.04), Inches(bar_max_w - 2.3), Inches(0.22),
        )
        track.fill.solid()
        track.fill.fore_color.rgb = BAR_BG
        track.line.fill.background()
        fill = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(left + 3.35), Inches(y + 0.04),
            Inches(max((bar_max_w - 2.3) * ratio, 0.08)), Inches(0.22),
        )
        fill.fill.solid()
        fill.fill.fore_color.rgb = CSC_RED
        fill.line.fill.background()
        detail = f"{num:,}/{denom:,}" if denom else f"{num:,}+"
        add_textbox(
            slide, left + bar_max_w - 0.85, y, 1.5, 0.28,
            [pct], size=10, bold=True, color=CSC_RED, align=PP_ALIGN.RIGHT,
        )
        add_textbox(slide, left + 3.35, y + 0.26, bar_max_w - 2.3, 0.22, [detail], size=8, color=TEXT_GRAY)
        y += 0.72


def parkwood_detail_rows():
    rows = []
    for r in PARKWOOD_ROOMS:
        monthly = r["monthly"]
        nine_mo = round(monthly * LEASE_MONTHS / ACADEMIC_MONTHS)
        area = round(r["unit_area"] / r["occupants"])
        psf = round(monthly / area)
        rows.append([r["label"], f"{monthly:,}", f"{nine_mo:,}", f"{area}（估）", f"{psf}（估）"])
    return rows


def slide_demand_gap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide)
    add_subsection_title(slide, "1、中大住宿需求缺口")

    half_w = 6.05
    gap = 0.53
    right_x = MARGIN + half_w + gap
    tbl_headers = ["指標", "數字", "備註"]
    tbl_h = 1.72

    add_table_banner(slide, MARGIN, 1.18, half_w, "本科住宿供需")
    add_data_table(
        slide, MARGIN, 1.58, half_w, tbl_h, tbl_headers, UG_DEMAND,
        col_widths=[1.85, 1.35, 2.85], header_size=10, body_size=9, left_cols={0, 2},
    )
    add_kpi_box(slide, MARGIN, 3.38, half_w, 0.95, "校内覆盖率", "52%", "9,500 / 18,438")

    add_table_banner(slide, right_x, 1.18, half_w, "研究生住宿供需")
    add_data_table(
        slide, right_x, 1.58, half_w, tbl_h, tbl_headers, PG_DEMAND,
        col_widths=[1.85, 1.35, 2.85], header_size=10, body_size=9, left_cols={0, 2},
    )
    add_kpi_box(slide, right_x, 3.38, half_w, 0.95, "PGH 覆盖率", "17%", "1,100 / 6,349")

    add_table_banner(slide, MARGIN, 4.48, CONTENT_W, "深港双城租住：香港入读 · 深圳居住")
    add_data_table(
        slide, MARGIN, 4.9, CONTENT_W, 1.75, ["维度", "数据/区间", "说明"], SHENZHEN_RENTAL,
        col_widths=[2.4, 2.1, 8.13], header_size=10, body_size=9, left_cols={0, 2},
    )

    add_textbox(
        slide, MARGIN, 6.78, CONTENT_W, 0.38,
        ["逻辑链：校内宿位不足 → 校外香港租务承压 → 部分港漂北上深圳（租金约为香港 1/3–1/2，但需每日跨境）"],
        size=10, color=CSC_RED, line_spacing=1.25,
    )


def slide_campus_detail(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide)
    add_subsection_title(slide, "1（續）、校內學生宿舍供給")

    add_table_banner(slide, MARGIN, 1.18, CONTENT_W, "校内宿位一览（九书院 + I-House / PGH / 在建，本科合计约 9,500+）")
    headers = ["项目", "宿位", "住宿模式", "费用参考", "备注"]
    hi = {i for i, row in enumerate(CAMPUS_ALL) if row[0] in FULL_RES_COLLEGES}
    add_data_table(
        slide, MARGIN, 1.58, CONTENT_W, 5.35, headers, CAMPUS_ALL,
        col_widths=[1.55, 0.85, 2.35, 1.75, 2.93],
        header_size=10, body_size=8, highlight_rows=hi, left_cols={0, 2, 4},
    )


def slide_offcampus_market(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide)
    add_subsection_title(slide, "2、校外學生公寓")

    add_table_banner(slide, MARGIN, 1.18, CONTENT_W, "市场化学生宿舍供应（近中大）— 学生公寓 · 长租酒店")

    add_data_table(
        slide, MARGIN, 1.58, CONTENT_W, 3.55, MARKET_DORM_HEADERS, MARKET_DORMS,
        col_widths=[0.5, 0.95, 2.15, 1.05, 1.55, 2.95, 0.85],
        header_size=9, body_size=8,
        left_cols={2, 5},
    )

    add_textbox(
        slide, MARGIN, 5.28, CONTENT_W, 0.72,
        [OFFCAMPUS_CONCLUSION],
        size=10, color=CSC_RED, line_spacing=1.3,
    )


def slide_parkwood_detail(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide)
    add_subsection_title(slide, "3、迎林 Parkwood — 房型定價深拆")

    summary = [
        "122伙改床位出租｜床位月租 5,335–5,950 元/月；预繳 11.5 個月 + 按金 11,000 元（HK01 2025/07）",
        "中大 2025 认可合作校外宿舍；2026 合作收尾中",
    ]
    add_textbox(slide, MARGIN, 1.18, CONTENT_W, 0.48, summary, size=10, color=TEXT_GRAY, line_spacing=1.25)

    add_table_banner(slide, MARGIN, 1.72, CONTENT_W, "房型定价拆解（租金 HK01；面积估算）")
    headers = [
        "房型",
        "月租金\n（元/月）",
        "折算到9个月\n月租（元/月）",
        "使用面积\n（平方呎）",
        "使用面积呎价\n（元/呎）",
    ]
    add_data_table(
        slide, MARGIN, 2.14, CONTENT_W, 1.75, headers, parkwood_detail_rows(),
        col_widths=[2.4, 2.0, 2.35, 2.0, 2.38],
        header_size=10, body_size=10,
        peach_rows=True, left_cols={0},
    )

    avg_unit = round(PARKWOOD_GFA_TOTAL / PARKWOOD_UNITS)
    note_lines = [
        f"折算9个月月租 = 月租金 × {LEASE_MONTHS} ÷ {ACADEMIC_MONTHS}",
        f"面积（估）= 户实用面积 ÷ 人数；户均参考批则 {PARKWOOD_GFA_TOTAL:,}呎÷{PARKWOOD_UNITS}户≈{avg_unit}呎",
        "首周申请 250+ 人（约可租床位 70%）；仍为近校学生公寓定价标竿",
    ]
    add_textbox(slide, MARGIN, 4.05, CONTENT_W, 0.75, note_lines, size=9, color=TEXT_GRAY, line_spacing=1.35)


def slide_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide)
    add_subsection_title(slide, "4、小結")

    add_table_banner(slide, MARGIN, 1.18, CONTENT_W, "投资判断摘要")
    headers = ["維度", "判斷"]
    rows = [
        ["本科需求", "校内覆盖约 52%，校外 9,000+ 人次/学年粗算缺口"],
        ["研究生需求", "PGH 覆盖约 17%；修课式/兼读/自资无校内宿"],
        ["深圳替代", "未获宿港漂北上；罗湖公寓约 40% 租客为在港读研内地生"],
        ["校内供给", "9,500+ 书院宿 + I-House/PGH；蔡继有（250）、E计划（+300）在建"],
        ["校外公寓", "近中大商业化宿舍床位仅数百量级，缺口极大；酒店长租 8,550–11,880/床"],
        ["对黄宜坳", "大埔走廊具填补缺口潜力；需实测通勤 + 对标迎林房型/呎租"],
    ]
    add_data_table(
        slide, MARGIN, 1.58, CONTENT_W, 4.85, headers, rows,
        col_widths=[1.45, 11.18], header_size=10, body_size=10,
        left_cols={0, 1},
    )


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_demand_gap(prs)
    slide_campus_detail(prs)
    slide_offcampus_market(prs)
    slide_parkwood_detail(prs)
    slide_summary(prs)

    try:
        prs.save(str(OUT_PATH))
        print(f"Saved: {OUT_PATH}")
    except PermissionError:
        prs.save(str(OUT_PATH_FALLBACK))
        print(f"Saved (fallback): {OUT_PATH_FALLBACK}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
