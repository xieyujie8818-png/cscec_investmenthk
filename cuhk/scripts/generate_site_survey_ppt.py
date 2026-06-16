# -*- coding: utf-8 -*-
"""Generate site survey PowerPoint for 大埔黃宜坳 project."""
from __future__ import annotations

import io
import math
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
PHOTOS_ROOT = ROOT / "項目實地圖片"
MAP_PATH = ROOT / "項目地塊地圖.png"
OUT_PATH = ROOT / "項目地塊實地周邊情況.pptx"

FONT_NAME = "Microsoft YaHei"  # 微軟雅黑
FONT_BASE = 15
FONT_TOP = 11  # 考察要點、情況說明
FONT_TITLE = 20
FONT_SECTION = 18
FONT_CAPTION = 13
FONT_SMALL = 12

CSC_RED = RGBColor(192, 0, 0)
TEXT_DARK = RGBColor(51, 51, 51)
TEXT_GRAY = RGBColor(102, 102, 102)
WHITE = RGBColor(255, 255, 255)
PANEL_BG = RGBColor(248, 248, 248)
FRAME_LINE = RGBColor(210, 210, 210)

# Content vertical bands (inches)
CONTENT_TOP = 1.08
CONTENT_MARGIN = 0.32
CONTENT_WIDTH = 12.7
TOP_BAND_H = 1.05
MAIN_TOP = CONTENT_TOP + TOP_BAND_H + 0.06
MAIN_BOTTOM = 7.18
COLUMN_GAP = 0.08
MAP_RATIO = 0.40
PHOTO_RATIO = 0.60
MAP_W = CONTENT_WIDTH * MAP_RATIO
PHOTO_W = CONTENT_WIDTH * PHOTO_RATIO
MAP_LEFT = CONTENT_MARGIN
PHOTO_LEFT = CONTENT_MARGIN + MAP_W + COLUMN_GAP

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".webp", ".heic"}


def list_images(folder: Path) -> list[Path]:
    if not folder.is_dir():
        return []
    files = [p for p in folder.iterdir() if p.suffix.lower() in IMAGE_EXTS and p.is_file()]
    return sorted(files, key=lambda p: p.name)


def set_run_font(run, size: int, bold: bool = False, color: RGBColor = TEXT_DARK):
    run.font.name = FONT_NAME
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    lines: list[str],
    size: int = 14,
    bold: bool = False,
    color: RGBColor = TEXT_DARK,
    align=PP_ALIGN.LEFT,
    line_spacing: float = 1.35,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.line_spacing = line_spacing
        if p.runs:
            set_run_font(p.runs[0], size, bold, color)
    return box


def add_header(slide, chapter: str = "二、項目概況"):
    """Top header with logo placeholder and chapter title."""
    # Logo placeholder (user adds later)
    placeholder = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.35), Inches(0.18), Inches(1.6), Inches(0.45)
    )
    placeholder.fill.background()
    placeholder.line.fill.background()
    add_textbox(
        slide,
        Inches(0.35),
        Inches(0.22),
        Inches(1.6),
        Inches(0.4),
        ["【Logo】"],
        size=10,
        color=TEXT_GRAY,
        align=PP_ALIGN.CENTER,
    )

    add_textbox(
        slide,
        Inches(8.8),
        Inches(0.15),
        Inches(4.2),
        Inches(0.55),
        [chapter],
        size=26,
        bold=True,
        color=CSC_RED,
        align=PP_ALIGN.RIGHT,
    )

    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.35), Inches(0.72), Inches(12.6), Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = CSC_RED
    line.line.fill.background()


def add_subsection_title(slide, title: str, top=Inches(0.8)):
    square = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.35), top, Inches(0.14), Inches(0.26)
    )
    square.fill.solid()
    square.fill.fore_color.rgb = CSC_RED
    square.line.fill.background()

    add_textbox(
        slide,
        Inches(0.56),
        top - Inches(0.03),
        Inches(12.2),
        Inches(0.42),
        [title],
        size=FONT_SECTION,
        bold=True,
        color=CSC_RED,
    )


def add_panel_header(slide, left, top, width, text: str, font_size: int = FONT_CAPTION):
    h = Inches(0.34)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = CSC_RED
    bar.line.fill.background()
    tf = bar.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    set_run_font(p.runs[0], font_size, True, WHITE)


def fit_image(left, top, max_w, max_h, img_path: Path):
    with Image.open(img_path) as im:
        w_px, h_px = im.size
    aspect = w_px / h_px
    max_w_f = max_w.inches if hasattr(max_w, "inches") else float(max_w)
    max_h_f = max_h.inches if hasattr(max_h, "inches") else float(max_h)
    w = max_w_f
    h = w / aspect
    if h > max_h_f:
        h = max_h_f
        w = h * aspect
    left_f = left.inches if hasattr(left, "inches") else float(left)
    top_f = top.inches if hasattr(top, "inches") else float(top)
    # center in box
    return Inches(left_f + (max_w_f - w) / 2), Inches(top_f + (max_h_f - h) / 2), Inches(w), Inches(h)


def _picture_stream(img_path: Path) -> io.BytesIO:
    """Normalize images to in-memory JPEG (handles MPO / HEIC from mobile cameras)."""
    with Image.open(img_path) as im:
        buf = io.BytesIO()
        im.convert("RGB").save(buf, format="JPEG", quality=92)
        buf.seek(0)
        return buf


def add_image(slide, img_path: Path, left, top, max_w, max_h):
    l, t, w, h = fit_image(left, top, max_w, max_h, img_path)
    source = _picture_stream(img_path)
    return slide.shapes.add_picture(source, l, t, width=w, height=h)


def add_photo_label(slide, text: str, left, top, width=Inches(1.35)):
    label_h = Inches(0.28)
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, label_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CSC_RED
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    set_run_font(p.runs[0], 11, True, WHITE)


def add_red_arrow(slide, x1, y1, x2, y2, width: float = 2.0):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = CSC_RED
    conn.line.width = Pt(width)
    conn.line.end_arrowhead = True
    return conn


def add_frame(slide, left, top, width, height, line_color=RGBColor(220, 220, 220)):
    frame = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height
    )
    frame.fill.background()
    frame.line.color.rgb = line_color
    frame.line.width = Pt(0.75)
    return frame


def add_map_pin(slide, cx, cy, label: str, view_dir: tuple[float, float] = (0.0, 1.0)):
    """
    Draggable location marker on map (separate shapes, not grouped).
    view_dir: normalized viewing direction for dashed sight-line arrow.
    """
    # Outer ring — user can drag to adjust
    ring_r = Inches(0.2)
    ring = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, cx - ring_r, cy - ring_r, ring_r * 2, ring_r * 2
    )
    ring.name = f"pin_{label}_ring"
    ring.fill.background()
    ring.line.color.rgb = CSC_RED
    ring.line.width = Pt(2.5)

    dot_r = Inches(0.09)
    dot = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, cx - dot_r, cy - dot_r, dot_r * 2, dot_r * 2
    )
    dot.name = f"pin_{label}_dot"
    dot.fill.solid()
    dot.fill.fore_color.rgb = CSC_RED
    dot.line.color.rgb = WHITE
    dot.line.width = Pt(1.5)

    # Label badge above pin
    badge_w, badge_h = Inches(0.42), Inches(0.3)
    badge = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        cx - badge_w / 2,
        cy - Inches(0.48),
        badge_w,
        badge_h,
    )
    badge.name = f"pin_{label}_label"
    badge.fill.solid()
    badge.fill.fore_color.rgb = CSC_RED
    badge.line.fill.background()
    tf = badge.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = label
    p.alignment = PP_ALIGN.CENTER
    set_run_font(p.runs[0], FONT_CAPTION, True, WHITE)

    # Viewing direction (dashed) — user can drag endpoints
    vx, vy = view_dir
    mag = math.hypot(vx, vy) or 1.0
    sight_len = Inches(0.38)
    ex = cx + sight_len * (vx / mag)
    ey = cy + sight_len * (vy / mag)
    sight = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, cx, cy, ex, ey)
    sight.name = f"pin_{label}_sight"
    sight.line.color.rgb = CSC_RED
    sight.line.width = Pt(1.75)
    sight.line.dash_style = 2
    sight.line.end_arrowhead = True

    return dot


def add_summary_top_band(slide, summary: str, page_idx: int = 0):
    """Top band: 情況說明 only, 11pt."""
    top = Inches(CONTENT_TOP)
    band_h = Inches(TOP_BAND_H)
    box_w = Inches(CONTENT_WIDTH)
    box_x = Inches(CONTENT_MARGIN)
    add_frame(slide, box_x, top, box_w, band_h, FRAME_LINE)
    add_panel_header(slide, box_x, top, box_w, "情況說明", FONT_TOP)
    text = summary if page_idx == 0 else f"{summary}（續頁補充實景。）"
    add_textbox(
        slide,
        box_x + Inches(0.14),
        top + Inches(0.4),
        box_w - Inches(0.28),
        band_h - Inches(0.46),
        [text],
        size=FONT_TOP,
        line_spacing=1.28,
    )


def add_map_panel_main(
    slide,
    map_path: Path,
    box_left,
    box_top,
    box_w,
    box_h,
    pin_xy: tuple[float, float],
    pin_label: str,
    view_dir: tuple[float, float],
    arrow_targets: list[tuple[int, int]],
) -> tuple[int, int]:
    """40% width map panel — larger image for readability."""
    add_frame(slide, box_left, box_top, box_w, box_h, FRAME_LINE)
    add_panel_header(slide, box_left, box_top, box_w, "考察節點位置示意（項目地塊地圖）")

    pad_in = 0.05
    img_top = box_top.inches + 0.36
    img_h = box_h.inches - 0.48
    img_left = box_left.inches + pad_in
    img_w = box_w.inches - pad_in * 2
    add_image(slide, map_path, Inches(img_left), Inches(img_top), Inches(img_w), Inches(img_h))

    px = Inches(img_left + img_w * pin_xy[0])
    py = Inches(img_top + img_h * pin_xy[1])
    add_map_pin(slide, px, py, pin_label, view_dir)

    add_textbox(
        slide,
        Inches(img_left),
        Inches(box_top.inches + box_h.inches - 0.18),
        Inches(img_w),
        Inches(0.16),
        ["紅色標記可拖曳微調坐標"],
        size=10,
        color=TEXT_GRAY,
        align=PP_ALIGN.CENTER,
    )

    for i, (tx, ty) in enumerate(arrow_targets):
        arr = add_red_arrow(slide, px, py, tx, ty, width=2.0)
        arr.name = f"pin_{pin_label}_to_photo_{i + 1}"

    return px, py


# Hub slide: pin on map (normalized) per node A–D
HUB_NODES = [
    {"label": "主入口", "pin_xy": (0.55, 0.14), "pin": "A", "view": (0.15, 1.0)},
    {"label": "入口內進入點", "pin_xy": (0.50, 0.30), "pin": "B", "view": (0.0, 1.0)},
    {"label": "第一階段地塊周邊", "pin_xy": (0.78, 0.50), "pin": "C", "view": (1.0, 0.2)},
    {"label": "地塊鄰近入口", "pin_xy": (0.60, 0.42), "pin": "D", "view": (-0.6, 0.8)},
]


def _edge_point_toward(shape, target_x, target_y):
    """Point on shape edge closest to target (for arrow endpoints)."""
    cx = shape.left + shape.width // 2
    cy = shape.top + shape.height // 2
    dx = target_x - cx
    dy = target_y - cy
    if dx == 0 and dy == 0:
        return cx, cy
    hw = shape.width // 2
    hh = shape.height // 2
    scale = min(abs(hw / dx) if dx else float("inf"), abs(hh / dy) if dy else float("inf"))
    return int(cx + dx * scale * 0.95), int(cy + dy * scale * 0.95)


def add_hub_slide(prs, map_path: Path, nodes: list[tuple[Path, str]]):
    """Overview: intro + central map + 4 corner photos + pins + arrows map→photo."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide)
    add_subsection_title(slide, "1、實地周邊環境")

    intro = (
        "項目地塊以鄉郊道路及規劃新建道路銜接周邊屋苑。實地考察沿主入口進入，依次觀察道路銜接、"
        "地塊邊界及一、二期發展範圍周邊地貌。地圖標記 A–D 對應右頁各考察節點實景。"
    )
    add_textbox(
        slide,
        Inches(CONTENT_MARGIN),
        Inches(1.15),
        Inches(CONTENT_WIDTH),
        Inches(0.72),
        [intro],
        size=FONT_TOP,
        line_spacing=1.3,
    )

    vis_top = Inches(1.92)
    vis_h = Inches(MAIN_BOTTOM - 1.92)
    vis_w = Inches(CONTENT_WIDTH)
    vis_left = Inches(CONTENT_MARGIN)
    add_frame(slide, vis_left, vis_top, vis_w, vis_h, FRAME_LINE)
    add_panel_header(slide, vis_left, vis_top, vis_w, "實地考察節點總覽")

    # Central map
    map_box_w = 4.85
    map_box_h = vis_h.inches - 0.5
    map_left = CONTENT_MARGIN + (CONTENT_WIDTH - map_box_w) / 2
    map_top = 1.92 + 0.42
    map_shape = add_image(
        slide,
        map_path,
        Inches(map_left),
        Inches(map_top),
        Inches(map_box_w),
        Inches(map_box_h),
    )

    # Four corner photos (larger)
    photo_slots = [
        (CONTENT_MARGIN + 0.06, 1.92 + 0.38, 3.55, (vis_h.inches - 0.45) / 2 - 0.04),  # NW
        (CONTENT_MARGIN + CONTENT_WIDTH - 3.61, 1.92 + 0.38, 3.55, (vis_h.inches - 0.45) / 2 - 0.04),  # NE
        (CONTENT_MARGIN + 0.06, 1.92 + 0.38 + (vis_h.inches - 0.45) / 2 + 0.02, 3.55, (vis_h.inches - 0.45) / 2 - 0.04),  # SW
        (CONTENT_MARGIN + CONTENT_WIDTH - 3.61, 1.92 + 0.38 + (vis_h.inches - 0.45) / 2 + 0.02, 3.55, (vis_h.inches - 0.45) / 2 - 0.04),  # SE
    ]

    photo_shapes = []
    pin_centers = []

    for i, (img_path, label_text) in enumerate(nodes):
        sl = photo_slots[i]
        pic = add_image(slide, img_path, Inches(sl[0]), Inches(sl[1]), Inches(sl[2]), Inches(sl[3]))
        photo_shapes.append(pic)
        add_photo_label(slide, HUB_NODES[i]["label"], pic.left, pic.top - Inches(0.3), Inches(1.55))

        meta = HUB_NODES[i]
        px = map_shape.left + int(map_shape.width * meta["pin_xy"][0])
        py = map_shape.top + int(map_shape.height * meta["pin_xy"][1])
        add_map_pin(slide, px, py, meta["pin"], meta["view"])
        pin_centers.append((px, py))

    # Arrows: map pin → photo (drawn last, on top)
    for i, (pic, (px, py)) in enumerate(zip(photo_shapes, pin_centers)):
        tx, ty = _edge_point_toward(pic, px, py)
        arr = add_red_arrow(slide, px, py, tx, ty, width=2.5)
        arr.name = f"hub_arrow_{HUB_NODES[i]['pin']}"


def _photo_grid_layout(n: int) -> list[tuple[float, float, float, float]]:
    """Horizontal-first grid — maximises photo size for leadership review."""
    gap = 0.012
    if n == 1:
        return [(0, 0, 1, 1)]
    if n == 2:
        w = (1 - gap) / 2
        return [(0, 0, w, 1), (w + gap, 0, w, 1)]
    if n == 3:
        w = (1 - gap * 2) / 3
        return [
            (0, 0, w, 1),
            (w + gap, 0, w, 1),
            (2 * (w + gap), 0, w, 1),
        ]
    w = (1 - gap) / 2
    h = (1 - gap) / 2
    return [
        (0, 0, w, h),
        (w + gap, 0, w, h),
        (0, h + gap, w, h),
        (w + gap, h + gap, w, h),
    ]


def add_detail_slides(
    prs,
    map_path: Path,
    subtitle: str,
    node_intro: str,
    bullets: list[str],
    images: list[Path],
    pin_xy: tuple[float, float],
    pin_label: str,
    view_dir: tuple[float, float],
    summary: str,
    per_page: int = 4,
):
    pages = max(1, math.ceil(len(images) / per_page))
    main_h = MAIN_BOTTOM - MAIN_TOP

    for page_idx in range(pages):
        chunk = images[page_idx * per_page : (page_idx + 1) * per_page]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_header(slide)
        title = f"2、實地節點實景 — {subtitle}"
        if page_idx > 0:
            title += "（續）"
        add_subsection_title(slide, title)

        add_summary_top_band(slide, summary, page_idx)

        # ── Bottom left: map 40% ──
        map_box_top = Inches(MAIN_TOP)
        map_box_h = Inches(main_h)
        map_box_left = Inches(MAP_LEFT)
        map_box_w = Inches(MAP_W)

        # ── Bottom right: photos 60% ──
        grid_left = Inches(PHOTO_LEFT)
        grid_top = Inches(MAIN_TOP)
        grid_w = Inches(PHOTO_W)
        grid_h = Inches(main_h)
        add_frame(slide, grid_left, grid_top, grid_w, grid_h, FRAME_LINE)
        add_panel_header(
            slide, grid_left, grid_top, grid_w, f"實地照片（本節點共 {len(images)} 張）"
        )

        n = len(chunk)
        positions = _photo_grid_layout(n)
        gw = PHOTO_W - 0.12
        gh = main_h - 0.42
        gl = PHOTO_LEFT + 0.06
        gt = MAIN_TOP + 0.4
        cell_gap = 0.05
        photo_shapes = []

        for idx, (img_path, (cx, cy, cw, ch)) in enumerate(zip(chunk, positions[:n])):
            cell_w = gw * cw - cell_gap
            cell_h = gh * ch - cell_gap
            cell_l = Inches(gl + gw * cx + (cell_gap / 2 if cx > 0 else 0))
            cell_t = Inches(gt + gh * cy + (cell_gap / 2 if cy > 0 else 0))
            pic = add_image(slide, img_path, cell_l, cell_t, Inches(cell_w), Inches(cell_h))
            photo_shapes.append(pic)

            badge = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.OVAL,
                cell_l + Inches(0.08),
                cell_t + Inches(0.08),
                Inches(0.32),
                Inches(0.32),
            )
            badge.fill.solid()
            badge.fill.fore_color.rgb = CSC_RED
            badge.line.fill.background()
            tf = badge.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.text = str(idx + 1 + page_idx * per_page)
            p.alignment = PP_ALIGN.CENTER
            set_run_font(p.runs[0], FONT_CAPTION, True, WHITE)

        arrow_targets = [(pic.left, pic.top + pic.height // 2) for pic in photo_shapes]
        add_map_panel_main(
            slide,
            map_path,
            map_box_left,
            map_box_top,
            map_box_w,
            map_box_h,
            pin_xy,
            pin_label,
            view_dir,
            arrow_targets,
        )


def add_summary_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide)
    add_subsection_title(slide, "實地考察小結")

    conclusions = [
        "1. 區位與供給環境：地塊處黃宜坳山坳腹地，一、二期依山形分布並以新建道路串聯；"
        "周邊為成熟中低密度住宅帶，區內新增供應有限，有利一期高尚住宅定位。",
        "2. 地塊物理條件：實地植被覆蓋良好，一期東側地塊邊界清晰；坡度及地貌與規劃中密度住宅"
        "開發強度大致匹配，未見明顯不利因素。",
        "3. 開發可達性：主入口及地塊鄰近入口現具基本車行條件，可支撐前期勘測及施工籌備；"
        "新建道路落成後可達性將進一步改善，需納入開發期及成本測算。",
        "4. 鄰里與變現參考：毗鄰雍怡雅苑、大埔寶馬山等可比屋苑，具成交對標基礎；"
        "自然護理區及海景資源強化產品溢價敘事，惟最終指標以政府審批為準。",
    ]
    add_textbox(
        slide,
        Inches(0.5),
        Inches(1.3),
        Inches(12.3),
        Inches(5.5),
        conclusions,
        size=14,
        line_spacing=1.55,
    )
    add_textbox(
        slide,
        Inches(0.5),
        Inches(6.85),
        Inches(12.3),
        Inches(0.35),
        ["以上為實地考察匯報，規劃及財務指標以政府審批及正式盡調為準"],
        size=10,
        color=TEXT_GRAY,
        align=PP_ALIGN.CENTER,
    )


def main():
    if not MAP_PATH.is_file():
        raise FileNotFoundError(f"Map not found: {MAP_PATH}")

    node_folders = [
        ("入口", "主入口"),
        ("入口進去點", "入口內進入點"),
        ("第一塊地塊附近", "第一階段地塊周邊"),
        ("第一塊地塊附近入口", "地塊鄰近入口"),
    ]
    node_images: list[tuple[Path, str]] = []
    for folder_name, label in node_folders:
        imgs = list_images(PHOTOS_ROOT / folder_name)
        if not imgs:
            raise FileNotFoundError(f"No images in {folder_name}")
        node_images.append((imgs[0], label))

    # pin_xy: normalized position within map image (x, y from 0–1)
    # pin_xy / view_dir: normalized; adjust in PPT by dragging shapes named pin_X_*
    detail_sections = [
        {
            "subtitle": "主入口",
            "node_intro": "本節點為項目外部主入口，銜接規劃新建道路北段與周邊鄉郊路網。",
            "bullets": [
                "• 現況道路等級、路面寬度及彎道坡度",
                "• 與周邊鄉郊路網／屋苑支路的銜接關係",
                "• 施工期車輛進出可行性及臨時便道需求初判",
            ],
            "folder": "入口",
            "pin_xy": (0.55, 0.14),
            "pin_label": "A",
            "view_dir": (0.15, 1.0),
            "summary": (
                "主入口位於規劃新建道路北段起點一帶，現以鄉郊小路與外部路網銜接，路面可容小型車輛通行；"
                "彎道坡度整體適中，具備前期勘測及施工車輛進出基礎，惟重型車輛及高峰期交通組織需於施工方案中專項評估。"
            ),
        },
        {
            "subtitle": "入口內進入點",
            "node_intro": "本節點為進入地盤後的首個關鍵觀測位置，用於核實道路對位及視線條件。",
            "bullets": [
                "• 進入地盤後視線、安全距離及兩側地貌（圍欄／植被／既有構築物）",
                "• 與規劃新建道路走向的對位關係",
            ],
            "folder": "入口進去點",
            "pin_xy": (0.50, 0.30),
            "pin_label": "B",
            "view_dir": (0.0, 1.0),
            "summary": (
                "入口內進入點處於地盤外圍過渡帶，視野受兩側山林圍合，實地走向與規劃新建道路大致吻合；"
                "兩側植被茂密、邊界清晰，清表及邊坡工程階段需同步考慮水土保持與施工安全距離控制。"
            ),
        },
        {
            "subtitle": "第一階段地塊周邊",
            "node_intro": "本節點位於第一階段（東側）高尚住宅發展範圍周邊，重點觀察地塊邊界及景觀條件。",
            "bullets": [
                "• 一期（東側）地塊邊界地貌：坡度、植被覆蓋、開挖難度初判",
                "• 與毗鄰屋苑（雍怡雅苑、大埔寶馬山等方向）的空間關係",
                "• 遠景視野（吐露港／船灣淡水湖方向，若照片可見）",
            ],
            "folder": "第一塊地塊附近",
            "pin_xy": (0.78, 0.50),
            "pin_label": "C",
            "view_dir": (1.0, 0.2),
            "summary": (
                "第一階段（東側）地塊周邊林木覆蓋良好，邊界地貌以緩坡為主，開挖難度整體可控；"
                "與毗鄰屋苑保持適當距離，遠眺方向具吐露港及船灣淡水湖景觀資源，有利高尚住宅產品定位及溢價敘事。"
            ),
        },
        {
            "subtitle": "地塊鄰近入口",
            "node_intro": "本節點為一、二期地塊串聯的次級銜接位置，用於驗證分期開發動線。",
            "bullets": [
                "• 次級入口／銜接點現況及與二期（西側）地塊的空間關係",
                "• 一、二期地塊串聯動線的實地驗證",
            ],
            "folder": "第一塊地塊附近入口",
            "pin_xy": (0.60, 0.42),
            "pin_label": "D",
            "view_dir": (-0.6, 0.8),
            "summary": (
                "地塊鄰近入口為一、二期地塊串聯的次級銜接節點，現況以林間小徑及台階為主，與二期（西側）地塊高差及植被條件相若；"
                "整體動線可透過新建道路統籌貫通，施工期需重點評估便道設置及與主入口的交通分流安排。"
            ),
        },
    ]

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    add_hub_slide(prs, MAP_PATH, node_images)
    for sec in detail_sections:
        add_detail_slides(
            prs,
            MAP_PATH,
            sec["subtitle"],
            sec["node_intro"],
            sec["bullets"],
            list_images(PHOTOS_ROOT / sec["folder"]),
            sec["pin_xy"],
            sec["pin_label"],
            sec["view_dir"],
            sec["summary"],
        )
    add_summary_slide(prs)

    prs.save(str(OUT_PATH))
    print(f"Saved: {OUT_PATH}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
